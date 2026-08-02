from pathlib import Path

import pytest

from src.utils.document_text import document_to_markdown


def test_document_to_markdown_reads_supported_office_formats(tmp_path: Path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from reportlab.pdfgen.canvas import Canvas

    docx_path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("DOCX content")
    document.save(docx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active.append(["XLSX content", 7])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PPTX content"
    presentation.save(pptx_path)

    pdf_path = tmp_path / "sample.pdf"
    canvas = Canvas(str(pdf_path))
    canvas.drawString(72, 720, "PDF content")
    canvas.save()

    assert "DOCX content" in document_to_markdown(docx_path)
    assert "XLSX content" in document_to_markdown(xlsx_path)
    assert "PPTX content" in document_to_markdown(pptx_path)
    assert "PDF content" in document_to_markdown(pdf_path)


def test_document_to_markdown_rejects_unsupported_format(tmp_path: Path) -> None:
    path = tmp_path / "document.txt"
    path.write_text("plain text", encoding="utf-8")

    with pytest.raises(ValueError, match="Unsupported document format"):
        document_to_markdown(path)
