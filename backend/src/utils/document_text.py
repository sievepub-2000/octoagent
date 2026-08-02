"""Extract Markdown-friendly text from the four supported document formats."""

from pathlib import Path


def document_to_markdown(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        from docx import Document

        document = Document(path)
        blocks = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            rows = [[cell.text.strip().replace("|", "\\|") for cell in row.cells] for row in table.rows]
            if rows:
                blocks.extend(("| " + " | ".join(row) + " |") for row in rows)
        return "\n\n".join(blocks)

    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(path, data_only=True, read_only=True)
        sheets: list[str] = []
        try:
            for sheet in workbook.worksheets:
                rows = [
                    ["" if value is None else str(value).replace("|", "\\|") for value in row]
                    for row in sheet.iter_rows(values_only=True)
                ]
                rows = [row for row in rows if any(cell for cell in row)]
                if rows:
                    sheets.append(f"## {sheet.title}\n\n" + "\n".join("| " + " | ".join(row) + " |" for row in rows))
        finally:
            workbook.close()
        return "\n\n".join(sheets)

    if suffix == ".pptx":
        from pptx import Presentation

        slides: list[str] = []
        for number, slide in enumerate(Presentation(path).slides, 1):
            text = [str(shape.text).strip() for shape in slide.shapes if hasattr(shape, "text") and str(shape.text).strip()]
            if text:
                slides.append(f"## Slide {number}\n\n" + "\n\n".join(text))
        return "\n\n".join(slides)

    if suffix == ".pdf":
        from pypdf import PdfReader

        pages = [f"## Page {number}\n\n{text.strip()}" for number, page in enumerate(PdfReader(path).pages, 1) if (text := (page.extract_text() or ""))]
        return "\n\n".join(pages)

    raise ValueError(f"Unsupported document format: {suffix or '<none>'}")
