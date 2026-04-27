#!/usr/bin/env python3
"""Deep simulation test for OctoAgent — end-to-end workflow testing.

Tests:
  1. Create workspaces (single, branch, group modes)
  2. Compile plans via Brain Core
  3. Send agent messages for paper writing tasks
  4. Generate docx and ppt outputs
  5. Verify all outputs
"""

import json
import os
import sys
import time
import subprocess
from pathlib import Path
from urllib.request import urlopen, Request
from urllib.error import URLError, HTTPError

GATEWAY = os.environ.get("GATEWAY_URL", "http://127.0.0.1:19882")
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "test_outputs"

# Three complex paper topics
PAPER_TOPICS = [
    {
        "title": "量子计算在密码学中的应用与挑战",
        "goal": (
            "Write a comprehensive research paper on 'Applications and Challenges of "
            "Quantum Computing in Cryptography'. Cover: 1) Quantum computing fundamentals "
            "(qubits, entanglement, superposition), 2) Shor's algorithm and RSA vulnerability, "
            "3) Post-quantum cryptography algorithms (lattice-based, hash-based, code-based), "
            "4) Current state of quantum hardware, 5) Timeline predictions and migration strategies. "
            "Output a well-structured markdown file with abstract, sections, and references."
        ),
        "filename": "quantum_crypto_paper",
    },
    {
        "title": "大语言模型在自动化科学发现中的角色",
        "goal": (
            "Write a research paper on 'The Role of Large Language Models in Automated Scientific "
            "Discovery'. Cover: 1) LLM capabilities for hypothesis generation, 2) AlphaFold and "
            "protein structure prediction, 3) LLM-driven drug discovery pipelines, 4) Automated "
            "theorem proving and mathematical reasoning, 5) Limitations: hallucination, "
            "reproducibility, and validation challenges. Include case studies and future outlook."
        ),
        "filename": "llm_scientific_discovery_paper",
    },
    {
        "title": "去中心化人工智能：联邦学习与区块链的融合",
        "goal": (
            "Write a research paper on 'Decentralized AI: Convergence of Federated Learning "
            "and Blockchain'. Cover: 1) Federated learning architecture and privacy guarantees, "
            "2) Blockchain-based model aggregation and incentive mechanisms, 3) Differential "
            "privacy integration, 4) Real-world deployments (healthcare, IoT, finance), "
            "5) Scalability challenges and solutions. Provide technical depth with diagrams "
            "described in text and comprehensive references."
        ),
        "filename": "decentralized_ai_paper",
    },
]

WORKFLOW_MODES = ["single", "branch", "group"]

# Map agent per test
AGENT_ASSIGNMENTS = {
    "single": "qwen-researcher",
    "branch": "nemotron-writer",
    "group": "qwen-researcher",
}


def api_call(method: str, path: str, data: dict | None = None, timeout: int = 120) -> dict | list | None:
    """Make an API call to the gateway."""
    url = f"{GATEWAY}{path}"
    headers = {"Content-Type": "application/json"}
    body = json.dumps(data).encode() if data else None
    req = Request(url, data=body, headers=headers, method=method)
    try:
        with urlopen(req, timeout=timeout) as resp:
            raw = resp.read().decode()
            return json.loads(raw) if raw else None
    except HTTPError as e:
        body_text = e.read().decode() if e.fp else ""
        print(f"  HTTP {e.code}: {body_text[:300]}")
        return {"error": True, "status": e.code, "detail": body_text}
    except URLError as e:
        print(f"  Connection error: {e.reason}")
        return {"error": True, "detail": str(e.reason)}
    except Exception as e:
        print(f"  Unexpected error: {e}")
        return {"error": True, "detail": str(e)}


def check_models():
    """Verify models are loaded."""
    print("\n=== Test: Model Configuration ===")
    result = api_call("GET", "/api/models")
    if not result or "error" in result:
        print("  FAIL: Cannot get models")
        return False
    models = result.get("models", [])
    print(f"  Models loaded: {len(models)}")
    for m in models:
        print(f"    - {m['name']:30s} fallback={m.get('fallback_models', [])}")
    if len(models) < 3:
        print("  FAIL: Expected at least 3 models")
        return False
    print("  PASS")
    return True


def check_agents():
    """Verify agents exist."""
    print("\n=== Test: Agent Configuration ===")
    result = api_call("GET", "/api/agents")
    if not result or "error" in result:
        print("  FAIL: Cannot get agents")
        return False
    agents = result.get("agents", [])
    print(f"  Agents: {len(agents)}")
    for a in agents:
        print(f"    - {a['name']:20s} model={a['model']}")
    agent_names = {a["name"] for a in agents}
    required = {"qwen-researcher", "nemotron-writer"}
    if not required.issubset(agent_names):
        print(f"  FAIL: Missing agents: {required - agent_names}")
        return False
    print("  PASS")
    return True


def create_workspace(mode: str, topic: dict) -> dict | None:
    """Create a task workspace."""
    print(f"\n  Creating {mode} workspace: {topic['title'][:40]}...")
    result = api_call("POST", "/api/task-workspaces", {
        "name": f"[{mode.upper()}] {topic['title']}",
        "goal": topic["goal"],
        "mode": mode,
        "summary": f"Deep simulation test — {mode} mode paper writing",
    })
    if not result or "error" in result:
        print(f"  FAIL: Create workspace error: {result}")
        return None
    task_id = result.get("task_id")
    agents = result.get("agents", [])
    print(f"  Created: task_id={task_id}, agents={len(agents)}, status={result.get('status')}")
    for a in agents:
        print(f"    Agent: {a['name']} (role={a['role']}, id={a['agent_id'][:20]}...)")
    return result


def compile_workspace(task_id: str) -> dict | None:
    """Compile workspace plan."""
    print(f"  Compiling workspace {task_id[:20]}...")
    result = api_call("POST", f"/api/task-workspaces/{task_id}/compile", timeout=180)
    if not result or "error" in result:
        print(f"  FAIL: Compile error: {result}")
        return None
    print(f"  Compiled: status={result.get('status')}")
    return result


def send_message(task_id: str, agent_id: str, content: str) -> dict | None:
    """Send a message to an agent."""
    print(f"  Sending message to agent {agent_id[:20]}...")
    result = api_call("POST", f"/api/task-workspaces/{task_id}/agents/{agent_id}/messages", {
        "content": content,
    }, timeout=300)
    if not result or "error" in result:
        print(f"  FAIL: Message error: {result}")
        return None
    messages = result.get("messages", [])
    assistant_msgs = [m for m in messages if m.get("role") == "assistant"]
    if assistant_msgs:
        last = assistant_msgs[-1]
        content_preview = (last.get("content") or "")[:200]
        print(f"  Response ({len(assistant_msgs)} assistant msgs): {content_preview}...")
    else:
        print(f"  WARNING: No assistant response in {len(messages)} messages")
    return result


def generate_documents(topic: dict, markdown_content: str):
    """Convert markdown to docx and pptx using Python."""
    OUTPUT_DIR.mkdir(exist_ok=True)
    base = topic["filename"]

    # Save markdown
    md_path = OUTPUT_DIR / f"{base}.md"
    md_path.write_text(markdown_content, encoding="utf-8")
    print(f"  Saved markdown: {md_path}")

    # Generate DOCX using pandoc
    docx_path = OUTPUT_DIR / f"{base}.docx"
    try:
        subprocess.run(
            ["pandoc", str(md_path), "-o", str(docx_path), "--from=markdown", "--to=docx"],
            check=True, capture_output=True, timeout=30,
        )
        print(f"  Generated DOCX: {docx_path} ({docx_path.stat().st_size} bytes)")
    except Exception as e:
        print(f"  DOCX generation failed: {e}")
        # Fallback: python-docx
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(topic["title"], level=0)
            for line in markdown_content.split("\n"):
                if line.startswith("# "):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith("## "):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith("### "):
                    doc.add_heading(line[4:], level=3)
                elif line.strip():
                    doc.add_paragraph(line)
            doc.save(str(docx_path))
            print(f"  Generated DOCX (fallback): {docx_path} ({docx_path.stat().st_size} bytes)")
        except Exception as e2:
            print(f"  DOCX fallback also failed: {e2}")

    # Generate PPTX
    pptx_path = OUTPUT_DIR / f"{base}.pptx"
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = topic["title"]
        if slide.placeholders[1]:
            slide.placeholders[1].text = "OctoAgent Deep Simulation Test"

        # Content slides from sections
        sections = []
        current_section = {"title": "", "bullets": []}
        for line in markdown_content.split("\n"):
            if line.startswith("## "):
                if current_section["title"]:
                    sections.append(current_section)
                current_section = {"title": line[3:].strip(), "bullets": []}
            elif line.startswith("### "):
                current_section["bullets"].append(line[4:].strip())
            elif line.strip() and not line.startswith("#"):
                # Take first sentence as bullet
                sent = line.strip()[:120]
                if sent and len(current_section["bullets"]) < 6:
                    current_section["bullets"].append(sent)
        if current_section["title"]:
            sections.append(current_section)

        for sec in sections[:10]:  # Max 10 content slides
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec["title"]
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(sec["bullets"][:6]):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet

        prs.save(str(pptx_path))
        print(f"  Generated PPTX: {pptx_path} ({pptx_path.stat().st_size} bytes)")
    except Exception as e:
        print(f"  PPTX generation failed: {e}")


def run_workflow_test(mode: str, topic: dict) -> bool:
    """Run a complete workflow test."""
    print(f"\n{'='*60}")
    print(f"  WORKFLOW TEST: {mode.upper()} mode — {topic['title']}")
    print(f"{'='*60}")

    # 1. Create workspace
    ws = create_workspace(mode, topic)
    if not ws:
        return False
    task_id = ws["task_id"]
    agents = ws.get("agents", [])

    # 2. Compile plan
    compiled = compile_workspace(task_id)
    if not compiled:
        print("  WARNING: Compile failed, proceeding with direct message...")

    # 3. Send task message to lead agent
    if not agents:
        print("  FAIL: No agents in workspace")
        return False

    lead_agent = agents[0]
    agent_id = lead_agent["agent_id"]

    # Construct detailed paper writing instruction
    prompt = (
        f"请用中英文混合方式撰写一篇关于「{topic['title']}」的学术论文。\n\n"
        f"具体要求：\n{topic['goal']}\n\n"
        f"论文结构必须包括：\n"
        f"1. Abstract（摘要）\n"
        f"2. Introduction（引言）\n"
        f"3. Background（背景知识）\n"
        f"4. Methodology / Technical Analysis（方法论/技术分析）\n"
        f"5. Case Studies / Applications（案例研究/应用）\n"
        f"6. Discussion（讨论）\n"
        f"7. Conclusion（结论）\n"
        f"8. References（参考文献，至少列出10条）\n\n"
        f"请直接输出完整的Markdown格式论文，不要有任何其它说明。"
    )

    result = send_message(task_id, agent_id, prompt)
    if not result:
        return False

    # Extract the paper content
    messages = result.get("messages", [])
    assistant_msgs = [m for m in messages if m.get("role") == "assistant"]
    if not assistant_msgs:
        print("  FAIL: No assistant response")
        return False

    paper_content = assistant_msgs[-1].get("content", "")
    if len(paper_content) < 200:
        print(f"  WARNING: Paper content too short ({len(paper_content)} chars)")
        # Send follow-up
        result2 = send_message(task_id, agent_id,
            "请继续完成论文的剩余部分。输出完整的Markdown格式。")
        if result2:
            msgs2 = result2.get("messages", [])
            asst2 = [m for m in msgs2 if m.get("role") == "assistant"]
            if asst2:
                paper_content += "\n\n" + (asst2[-1].get("content") or "")

    print(f"  Paper content length: {len(paper_content)} chars")

    # 4. Generate documents
    if len(paper_content) > 100:
        print("\n  Generating documents...")
        generate_documents(topic, paper_content)
    else:
        print("  SKIP: Paper content too short for document generation")

    # 5. Check workspace status
    ws_status = api_call("GET", f"/api/task-workspaces/{task_id}")
    if ws_status:
        print(f"  Final workspace status: {ws_status.get('status')}")
        progress = ws_status.get("progress", {})
        print(f"  Progress: {progress}")

    return len(paper_content) > 200


def main():
    print("="*60)
    print("  OctoAgent Deep Simulation Test")
    print(f"  Gateway: {GATEWAY}")
    print(f"  Output: {OUTPUT_DIR}")
    print("="*60)

    OUTPUT_DIR.mkdir(exist_ok=True)
    results = {}

    # Pre-flight checks
    if not check_models():
        print("\nABORT: Model configuration failed")
        sys.exit(1)
    if not check_agents():
        print("\nABORT: Agent configuration failed")
        sys.exit(1)

    # Run three workflow tests
    for i, (mode, topic) in enumerate(zip(WORKFLOW_MODES, PAPER_TOPICS)):
        try:
            ok = run_workflow_test(mode, topic)
            results[f"{mode}_{topic['filename']}"] = "PASS" if ok else "FAIL"
        except Exception as e:
            print(f"  ERROR: {e}")
            results[f"{mode}_{topic['filename']}"] = f"ERROR: {e}"

    # Summary
    print("\n" + "="*60)
    print("  TEST SUMMARY")
    print("="*60)
    for test, result in results.items():
        status = "✓" if result == "PASS" else "✗"
        print(f"  {status} {test}: {result}")

    # List generated files
    print(f"\n  Output directory: {OUTPUT_DIR}")
    if OUTPUT_DIR.exists():
        for f in sorted(OUTPUT_DIR.iterdir()):
            print(f"    {f.name:40s} {f.stat().st_size:>8,} bytes")

    passed = sum(1 for v in results.values() if v == "PASS")
    total = len(results)
    print(f"\n  Result: {passed}/{total} tests passed")
    return 0 if passed == total else 1


if __name__ == "__main__":
    sys.exit(main())
