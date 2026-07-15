#!/usr/bin/env python3
"""Generate a small, standards-compliant Office/PDF/Markdown artifact."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any


def _text(value: Any) -> str:
    return "" if value is None else str(value)


def _sections(spec: dict[str, Any]) -> list[dict[str, Any]]:
    value = spec.get("sections", [])
    return [item for item in value if isinstance(item, dict)] if isinstance(value, list) else []


def generate_docx(spec: dict[str, Any], output: Path) -> None:
    from docx import Document

    document = Document()
    document.core_properties.title = _text(spec.get("title"))
    if spec.get("title"):
        document.add_heading(_text(spec["title"]), level=0)
    for section in _sections(spec):
        if section.get("heading"):
            document.add_heading(_text(section["heading"]), level=1)
        for paragraph in section.get("paragraphs", []):
            document.add_paragraph(_text(paragraph))
        for bullet in section.get("bullets", []):
            document.add_paragraph(_text(bullet), style="List Bullet")
    document.save(output)


def generate_xlsx(spec: dict[str, Any], output: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = _text(spec.get("sheet_name") or "Sheet1")[:31]
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if isinstance(headers, list) and headers:
        sheet.append([_text(item) for item in headers])
        for cell in sheet[1]:
            cell.font = Font(bold=True)
    if isinstance(rows, list):
        for row in rows:
            if isinstance(row, list):
                sheet.append(row)
    sheet.freeze_panes = "A2" if headers else None
    workbook.save(output)


def generate_pptx(spec: dict[str, Any], output: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slides = spec.get("slides", [])
    if not isinstance(slides, list) or not slides:
        slides = [{"title": spec.get("title", "Presentation"), "bullets": []}]
    for index, item in enumerate(slides):
        item = item if isinstance(item, dict) else {"title": _text(item)}
        layout = presentation.slide_layouts[0 if index == 0 else 1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = _text(item.get("title"))
        body = item.get("subtitle") if index == 0 else item.get("bullets", [])
        if len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            if isinstance(body, list):
                frame.clear()
                for bullet_index, bullet in enumerate(body):
                    paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                    paragraph.text = _text(bullet)
            else:
                frame.text = _text(body)
    presentation.save(output)


def generate_pdf(spec: dict[str, Any], output: Path) -> None:
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer

    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CJKTitle", parent=styles["Title"], fontName="STSong-Light")
    heading_style = ParagraphStyle("CJKHeading", parent=styles["Heading1"], fontName="STSong-Light")
    body_style = ParagraphStyle("CJKBody", parent=styles["BodyText"], fontName="STSong-Light", alignment=TA_LEFT, leading=16)
    story: list[Any] = []
    if spec.get("title"):
        story.extend([Paragraph(_text(spec["title"]), title_style), Spacer(1, 12)])
    for section in _sections(spec):
        if section.get("heading"):
            story.append(Paragraph(_text(section["heading"]), heading_style))
        for paragraph in section.get("paragraphs", []):
            story.extend([Paragraph(_text(paragraph), body_style), Spacer(1, 6)])
        bullets = [ListItem(Paragraph(_text(item), body_style)) for item in section.get("bullets", [])]
        if bullets:
            story.append(ListFlowable(bullets, bulletType="bullet"))
    SimpleDocTemplate(str(output), title=_text(spec.get("title"))).build(story)


def generate_markdown(spec: dict[str, Any], output: Path) -> None:
    lines = [f"# {_text(spec.get('title'))}", ""] if spec.get("title") else []
    for section in _sections(spec):
        if section.get("heading"):
            lines.extend([f"## {_text(section['heading'])}", ""])
        for paragraph in section.get("paragraphs", []):
            lines.extend([_text(paragraph), ""])
        for bullet in section.get("bullets", []):
            lines.append(f"- {_text(bullet)}")
        lines.append("")
    output.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


GENERATORS = {"docx": generate_docx, "xlsx": generate_xlsx, "pptx": generate_pptx, "pdf": generate_pdf, "md": generate_markdown}


def main() -> None:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    parser = argparse.ArgumentParser()
    parser.add_argument("--format", choices=sorted(GENERATORS), required=True)
    parser.add_argument("--spec", type=Path, required=True)
    parser.add_argument("--output-file", type=Path, required=True)
    args = parser.parse_args()
    if args.output_file.suffix.lower() != f".{args.format}":
        parser.error("output extension must match --format")
    spec = json.loads(args.spec.read_text(encoding="utf-8"))
    if not isinstance(spec, dict):
        parser.error("spec must be a JSON object")
    args.output_file.parent.mkdir(parents=True, exist_ok=True)
    GENERATORS[args.format](spec, args.output_file)
    print(json.dumps({"ok": True, "format": args.format, "path": str(args.output_file.resolve()), "bytes": args.output_file.stat().st_size}, ensure_ascii=False))


if __name__ == "__main__":
    main()
